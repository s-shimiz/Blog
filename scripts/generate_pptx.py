"""
Generate Agent Builder vs Copilot Studio presentation from template.
Content derived from: qiita/20260217_agent-builder-vs-copilot-studio.md
Template: Template/Copilot_Template_Public.potx (converted)
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from lxml import etree
import copy

TEMPLATE = r'C:\work\Blog\Template\Copilot_Template_Public_converted.pptx'
OUTPUT = r'C:\work\Blog\20260217_agent-builder-vs-copilot-studio.pptx'


def clone_layout_shapes(slide):
    """Clone non-placeholder shapes from layout to slide via XML."""
    layout = slide.slide_layout
    spTree = slide.shapes._spTree
    for shape in layout.shapes:
        if not shape.is_placeholder:
            clone = copy.deepcopy(shape._element)
            spTree.append(clone)


def find_shape(slide, name):
    """Find a shape by name in a slide."""
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def set_text_keep_format(shape, text):
    """Set text of a shape while preserving the first run's formatting."""
    if shape is None:
        return
    tf = shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        first_run = tf.paragraphs[0].runs[0]
        font_props = {
            'bold': first_run.font.bold,
            'italic': first_run.font.italic,
            'size': first_run.font.size,
            'color_rgb': first_run.font.color.rgb if first_run.font.color and first_run.font.color.type is not None else None,
            'name': first_run.font.name,
        }
    else:
        font_props = None

    # Handle multi-line text
    lines = text.split('\n')
    # Clear existing paragraphs
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)

    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
            p.clear()
            run = p.add_run()
        else:
            p = tf.add_paragraph()
            run = p.add_run()

        run.text = line
        if font_props:
            if font_props['bold'] is not None:
                run.font.bold = font_props['bold']
            if font_props['italic'] is not None:
                run.font.italic = font_props['italic']
            if font_props['size'] is not None:
                run.font.size = font_props['size']
            if font_props['color_rgb'] is not None:
                run.font.color.rgb = font_props['color_rgb']
            if font_props['name'] is not None:
                run.font.name = font_props['name']


def set_placeholder_text(slide, idx, text):
    """Set text of a placeholder by index."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            set_text_keep_format(ph, text)
            return


def add_slide(prs, layout_idx, clone_shapes=True):
    """Add a slide using the specified layout, optionally cloning non-placeholder shapes."""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    if clone_shapes:
        clone_layout_shapes(slide)
    return slide


def build_presentation():
    prs = Presentation(TEMPLATE)

    # Remove pre-existing slides from template
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # =========================================================================
    # Slide 1: Title Slide (Layout 0)
    # =========================================================================
    slide = add_slide(prs, 0)
    set_placeholder_text(slide, 0,
        'Agent Builder と Copilot Studio\n違いと使い分け')
    set_placeholder_text(slide, 10,
        'Microsoft 365 Copilot ライセンスで作るエージェント')

    # =========================================================================
    # Slide 2: Agenda (Layout 1)
    # =========================================================================
    slide = add_slide(prs, 1)
    set_placeholder_text(slide, 0, 'アジェンダ')
    set_placeholder_text(slide, 10,
        '1. Agent Builder とは\n'
        '2. Copilot Studio とは\n'
        '3. 機能比較\n'
        '4. 判断フロー ― どちらを選ぶべきか\n'
        '5. ガバナンスの違い\n'
        '6. Agent Builder → Copilot Studio への移行\n'
        '7. ライセンス要件\n'
        '8. まとめ')

    # =========================================================================
    # Slide 3: Section - Agent Builder (Layout 5)
    # =========================================================================
    slide = add_slide(prs, 5)
    set_text_keep_format(find_shape(slide, 'section_label'), '01')
    set_text_keep_format(find_shape(slide, 'section_title'), 'Agent Builder とは')
    set_text_keep_format(find_shape(slide, 'section_desc'),
        'Microsoft 365 Copilot アプリ内に組み込まれたエージェント作成機能')

    # =========================================================================
    # Slide 4: Agent Builder details (Layout 19: アイコン付きリスト)
    # =========================================================================
    slide = add_slide(prs, 19)
    set_text_keep_format(find_shape(slide, 'title'), 'Agent Builder の特徴')
    set_text_keep_format(find_shape(slide, 'subtitle'),
        'コード不要で、自然言語だけでエージェントを構築')

    set_text_keep_format(find_shape(slide, 'icon0'), '🌐')
    set_text_keep_format(find_shape(slide, 'item_title0'), 'アクセスポイント')
    set_text_keep_format(find_shape(slide, 'item_desc0'),
        'microsoft365.com/chat\noffice.com/chat\nMicrosoft Teams')

    set_text_keep_format(find_shape(slide, 'icon1'), '👤')
    set_text_keep_format(find_shape(slide, 'item_title1'), '想定ユーザー')
    set_text_keep_format(find_shape(slide, 'item_desc1'),
        'インフォメーション ワーカー\n（IT 部門以外も含む）')

    set_text_keep_format(find_shape(slide, 'icon2'), '📄')
    set_text_keep_format(find_shape(slide, 'item_title2'), '主な用途')
    set_text_keep_format(find_shape(slide, 'item_desc2'),
        '組織ナレッジを活用した\n軽量 Q&A エージェント')

    # =========================================================================
    # Slide 5: Section - Copilot Studio (Layout 5)
    # =========================================================================
    slide = add_slide(prs, 5)
    set_text_keep_format(find_shape(slide, 'section_label'), '02')
    set_text_keep_format(find_shape(slide, 'section_title'), 'Copilot Studio とは')
    set_text_keep_format(find_shape(slide, 'section_desc'),
        'スタンドアロン Web ポータルで高度なエージェントを構築')

    # =========================================================================
    # Slide 6: Copilot Studio details (Layout 19)
    # =========================================================================
    slide = add_slide(prs, 19)
    set_text_keep_format(find_shape(slide, 'title'), 'Copilot Studio の特徴')
    set_text_keep_format(find_shape(slide, 'subtitle'),
        'copilotstudio.microsoft.com で利用できるプロ向けツール')

    set_text_keep_format(find_shape(slide, 'icon0'), '🔧')
    set_text_keep_format(find_shape(slide, 'item_title0'), '高度な機能')
    set_text_keep_format(find_shape(slide, 'item_desc0'),
        'マルチステップ ワークフロー\n外部 API・カスタム コネクタ\nAzure AI サービス統合')

    set_text_keep_format(find_shape(slide, 'icon1'), '🏢')
    set_text_keep_format(find_shape(slide, 'item_title1'), '展開スコープ')
    set_text_keep_format(find_shape(slide, 'item_desc1'),
        '部署全体・組織全体\n外部顧客向けにも展開可能')

    set_text_keep_format(find_shape(slide, 'icon2'), '⚙️')
    set_text_keep_format(find_shape(slide, 'item_title2'), 'エンタープライズ ALM')
    set_text_keep_format(find_shape(slide, 'item_desc2'),
        'バージョン管理\ndev/test/prod 環境分離\nRBAC・テレメトリ')

    # =========================================================================
    # Slide 7: Section - 機能比較 (Layout 5)
    # =========================================================================
    slide = add_slide(prs, 5)
    set_text_keep_format(find_shape(slide, 'section_label'), '03')
    set_text_keep_format(find_shape(slide, 'section_title'), '機能比較')
    set_text_keep_format(find_shape(slide, 'section_desc'),
        'Agent Builder と Copilot Studio を多角的に比較')

    # =========================================================================
    # Slide 8: Comparison - Overview (Layout 7: 2カラム比較)
    # =========================================================================
    slide = add_slide(prs, 7)
    set_text_keep_format(find_shape(slide, 'title'), '機能比較：概要')
    set_text_keep_format(find_shape(slide, 'subtitle'),
        '対象ユーザー・展開スコープ・主な機能の違い')

    set_text_keep_format(find_shape(slide, 'label_before'), 'Agent Builder')
    set_text_keep_format(find_shape(slide, 'col1_text'),
        '▸ アクセス: M365 Copilot アプリ内\n'
        '▸ 対象: 自分自身・小規模チーム\n'
        '▸ 自然言語オーサリング\n'
        '▸ Microsoft Graph ベース Q&A\n'
        '▸ M365 アクセス許可を尊重\n'
        '▸ 管理: M365 管理センター')

    set_text_keep_format(find_shape(slide, 'label_after'), 'Copilot Studio')
    set_text_keep_format(find_shape(slide, 'col2_text'),
        '▸ アクセス: copilotstudio.microsoft.com\n'
        '▸ 対象: 部署・組織・外部顧客\n'
        '▸ マルチステップ ワークフロー\n'
        '▸ 外部 API・カスタム コネクタ\n'
        '▸ Azure AI サービス統合\n'
        '▸ 管理: Power Platform 管理センター')

    # =========================================================================
    # Slide 9: Section - 判断フロー (Layout 5)
    # =========================================================================
    slide = add_slide(prs, 5)
    set_text_keep_format(find_shape(slide, 'section_label'), '04')
    set_text_keep_format(find_shape(slide, 'section_title'), '判断フロー')
    set_text_keep_format(find_shape(slide, 'section_desc'),
        'どちらを選ぶべきか ― 4 つの観点で検討')

    # =========================================================================
    # Slide 10: Decision criteria (Layout 17: 番号付きステップ横型)
    # =========================================================================
    slide = add_slide(prs, 17)
    set_text_keep_format(find_shape(slide, 'title'), 'ツール選択の 4 つの観点')
    set_text_keep_format(find_shape(slide, 'subtitle'),
        '次の観点から最適なツールを判断します')

    set_text_keep_format(find_shape(slide, 'label0'), '対象ユーザー')
    set_text_keep_format(find_shape(slide, 'desc0'),
        '誰がエージェントを\n使うか？')

    set_text_keep_format(find_shape(slide, 'label1'), '展開スコープ')
    set_text_keep_format(find_shape(slide, 'desc1'),
        'どの範囲に\n共有するか？')

    set_text_keep_format(find_shape(slide, 'label2'), '機能')
    set_text_keep_format(find_shape(slide, 'desc2'),
        'どんなタスクを\n実行させるか？')

    set_text_keep_format(find_shape(slide, 'label3'), 'ガバナンス')
    set_text_keep_format(find_shape(slide, 'desc3'),
        'ALM が\n必要か？')

    # =========================================================================
    # Slide 11: Agent Builder use cases (Layout 11: 3カラム アクセントカラー付き)
    # =========================================================================
    slide = add_slide(prs, 11)
    set_text_keep_format(find_shape(slide, 'title'), 'Agent Builder を選ぶケース')
    set_text_keep_format(find_shape(slide, 'subtitle'),
        'すばやく・コード不要で・小規模に始めたい場合')

    set_text_keep_format(find_shape(slide, 'label0'), 'プロジェクト\nFAQ ボット')
    set_text_keep_format(find_shape(slide, 'desc0'),
        'プロジェクト ドキュメントをもとに\nよくある質問に回答')

    set_text_keep_format(find_shape(slide, 'label1'), '製品ドキュメント\nアシスタント')
    set_text_keep_format(find_shape(slide, 'desc1'),
        '社内の製品マニュアルや\nWiki から情報を検索')

    set_text_keep_format(find_shape(slide, 'label2'), 'オンボーディング\nエージェント')
    set_text_keep_format(find_shape(slide, 'desc2'),
        '新メンバーが社内\nナレッジ ベースから回答を得る')

    # =========================================================================
    # Slide 12: Copilot Studio use cases (Layout 11: 3カラム)
    # =========================================================================
    slide = add_slide(prs, 11)
    set_text_keep_format(find_shape(slide, 'title'), 'Copilot Studio を選ぶケース')
    set_text_keep_format(find_shape(slide, 'subtitle'),
        '組織全体・外部向け・複雑なワークフローが必要な場合')

    set_text_keep_format(find_shape(slide, 'label0'), 'カスタマー\nサポート')
    set_text_keep_format(find_shape(slide, 'desc0'),
        'サポート チケットの作成\n人間へのエスカレーション')

    set_text_keep_format(find_shape(slide, 'label1'), 'IT ヘルプデスク\nトリアージ')
    set_text_keep_format(find_shape(slide, 'desc1'),
        'IT リクエストの受付\nサポート チームへの振り分け')

    set_text_keep_format(find_shape(slide, 'label2'), '営業アシスタント\n（CRM 連携）')
    set_text_keep_format(find_shape(slide, 'desc2'),
        '売上データ取得・メモ作成\n承認ワークフローの起動')

    # =========================================================================
    # Slide 13: Section - ガバナンス (Layout 5)
    # =========================================================================
    slide = add_slide(prs, 5)
    set_text_keep_format(find_shape(slide, 'section_label'), '05')
    set_text_keep_format(find_shape(slide, 'section_title'), 'ガバナンスの違い')
    set_text_keep_format(find_shape(slide, 'section_desc'),
        '管理・監査・コンプライアンスの観点で比較')

    # =========================================================================
    # Slide 14: Governance comparison (Layout 21: 強調パネル)
    # =========================================================================
    slide = add_slide(prs, 21)
    set_text_keep_format(find_shape(slide, 'title'), 'ガバナンス比較')
    set_text_keep_format(find_shape(slide, 'subtitle'),
        '管理ポイントとセキュリティ制御の違い')

    set_text_keep_format(find_shape(slide, 'panel_title0'),
        'Agent Builder のガバナンス')
    set_text_keep_format(find_shape(slide, 'panel_desc0'),
        '▸ 新しい特権なし ― 既存の M365 アクセス許可を尊重\n'
        '▸ 組み込みの監査 ― 標準監査ログ・DLP・保持ポリシー適用\n'
        '▸ M365 管理センターでエージェントの管理・制御')

    set_text_keep_format(find_shape(slide, 'panel_title1'),
        'Copilot Studio のガバナンス')
    set_text_keep_format(find_shape(slide, 'panel_desc1'),
        '▸ 構造化開発（ALM）― dev/test/prod 環境の分離\n'
        '▸ コネクタ ガバナンス ― 接続先を管理者が制御\n'
        '▸ 環境レベルの DLP・RBAC・監査ポリシー\n'
        '▸ 管理者承認による組織カタログへの発行制御\n'
        '▸ Power Platform 管理センターで一元管理')

    # =========================================================================
    # Slide 15: Section - 移行・ライセンス (Layout 5)
    # =========================================================================
    slide = add_slide(prs, 5)
    set_text_keep_format(find_shape(slide, 'section_label'), '06')
    set_text_keep_format(find_shape(slide, 'section_title'), '移行とライセンス')
    set_text_keep_format(find_shape(slide, 'section_desc'),
        'Agent Builder から Copilot Studio への移行パスとライセンス要件')

    # =========================================================================
    # Slide 16: Migration steps (Layout 16: 縦3つステップ)
    # =========================================================================
    slide = add_slide(prs, 16)
    set_text_keep_format(find_shape(slide, 'title'),
        'Agent Builder → Copilot Studio への移行')
    set_text_keep_format(find_shape(slide, 'subtitle'),
        'スモールスタートで始めて、必要に応じて拡張')

    set_text_keep_format(find_shape(slide, 'step_title0'),
        'Agent Builder で開始')
    set_text_keep_format(find_shape(slide, 'step_desc0'),
        '自然言語でエージェントを作成\nSharePoint・Web をナレッジに設定')

    set_text_keep_format(find_shape(slide, 'step_title1'),
        'Copilot Studio にコピー')
    set_text_keep_format(find_shape(slide, 'step_desc1'),
        'コア構成とインストラクションを\nそのまま引き継ぎ')

    set_text_keep_format(find_shape(slide, 'step_title2'),
        '高度な機能で拡張')
    set_text_keep_format(find_shape(slide, 'step_desc2'),
        'ALM・外部 API 統合・\n分析ダッシュボード・ガバナンスを追加')

    # =========================================================================
    # Slide 17: License (Layout 39: シンプルリスト+補足パネル)
    # =========================================================================
    slide = add_slide(prs, 39)
    set_text_keep_format(find_shape(slide, 'title'), 'ライセンス要件')
    set_text_keep_format(find_shape(slide, 'subtitle'),
        'Microsoft 365 Copilot ライセンスがあれば追加費用なし')

    set_text_keep_format(find_shape(slide, 'list'),
        '▸ M365 Copilot アドオン ライセンス\n'
        '  → Agent Builder ✅ / Copilot Studio ✅\n\n'
        '▸ Copilot Credits / 従量課金\n'
        '  → Agent Builder ✅ / Copilot Studio ✅\n\n'
        '▸ 無料（指示 + 公開 Web のナレッジのみ）\n'
        '  → Agent Builder ✅ / Copilot Studio ❌')

    set_text_keep_format(find_shape(slide, 'note_title'), '💡 ポイント')
    set_text_keep_format(find_shape(slide, 'note_body'),
        'Microsoft 365 Copilot ライセンスがあれば、\n'
        'どちらのツールも追加費用なしで利用できます。')

    # =========================================================================
    # Slide 18: Summary (Layout 40: 対比+結論)
    # =========================================================================
    slide = add_slide(prs, 40)
    set_text_keep_format(find_shape(slide, 'title'), 'まとめ')
    set_text_keep_format(find_shape(slide, 'subtitle'),
        '選択基準と推奨アプローチ')

    set_text_keep_format(find_shape(slide, 'labelA'), 'Agent Builder')
    set_text_keep_format(find_shape(slide, 'descA'),
        '⭐ すばやく Q&A ボットを作りたい\n'
        '⭐ コードなしで完結したい\n'
        '⭐ 小規模チーム向け')

    set_text_keep_format(find_shape(slide, 'labelB'), 'Copilot Studio')
    set_text_keep_format(find_shape(slide, 'descB'),
        '⭐ 組織全体・外部向けに展開\n'
        '⭐ 外部システム連携・ワークフロー\n'
        '⭐ エンタープライズ ガバナンス')

    set_text_keep_format(find_shape(slide, 'conclusion_label'), '💡 推奨アプローチ')
    set_text_keep_format(find_shape(slide, 'conclusion_text'),
        '迷ったら、まず Agent Builder で小さく始めましょう。\n'
        '要件が拡大した時点で Copilot Studio にコピーして拡張できます。')

    # =========================================================================
    # Slide 19: References (Layout 1: Agenda/content)
    # =========================================================================
    slide = add_slide(prs, 1)
    set_placeholder_text(slide, 0, '参考リンク')
    set_placeholder_text(slide, 10,
        '▸ エージェントを構築する Microsoft 365 Copilot と Copilot Studio を選択する\n'
        '  https://learn.microsoft.com/microsoft-365-copilot/extensibility/copilot-studio-experience\n\n'
        '▸ Agent Builder in Microsoft 365 Copilot\n'
        '  https://learn.microsoft.com/microsoft-365-copilot/extensibility/agent-builder\n\n'
        '▸ Copy an agent to Copilot Studio\n'
        '  https://learn.microsoft.com/microsoft-365-copilot/extensibility/copy-agent-to-copilot-studio\n\n'
        '▸ Copilot Studio でエージェントを構築する\n'
        '  https://learn.microsoft.com/microsoft-copilot-studio/microsoft-copilot-extend-copilot-extensions')

    # =========================================================================
    # Slide 20: Closing (Layout 3)
    # =========================================================================
    slide = add_slide(prs, 3)

    # Save
    prs.save(OUTPUT)
    print(f'Presentation saved to: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')


if __name__ == '__main__':
    build_presentation()
