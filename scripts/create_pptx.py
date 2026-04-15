"""
Markdown 記事 → PowerPoint 変換スクリプト
テンプレート: Template/Copilot_Template_Public.potx
"""

import zipfile
import shutil
import os
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


# === 設定 ===
TEMPLATE_PATH = r"c:\work\Blog\Template\Copilot_Template_Public.potx"
OUTPUT_PATH = r"c:\work\Blog\qiita\20260217_agent-builder-vs-copilot-studio.pptx"

# カラーパレット
COLOR_TITLE = RGBColor(0x00, 0x78, 0xD4)       # Microsoft Blue
COLOR_HEADING = RGBColor(0x24, 0x29, 0x2E)      # Dark gray
COLOR_BODY = RGBColor(0x33, 0x33, 0x33)         # Body text
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_ACCENT = RGBColor(0x00, 0x78, 0xD4)       # Accent blue
COLOR_LIGHT_BG = RGBColor(0xF3, 0xF2, 0xF1)     # Light background
COLOR_TABLE_HEADER = RGBColor(0x00, 0x78, 0xD4)  # Table header bg
COLOR_TABLE_ALT = RGBColor(0xF0, 0xF6, 0xFC)     # Alternating row bg
COLOR_STAR = RGBColor(0xFF, 0xB9, 0x00)          # Star color


def convert_potx_to_pptx(potx_path):
    """potx テンプレートを pptx として読み込めるように変換"""
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    tmp.close()
    shutil.copy2(potx_path, tmp.name)
    with zipfile.ZipFile(tmp.name, "r") as zin:
        contents = {}
        for item in zin.namelist():
            data = zin.read(item)
            if item == "[Content_Types].xml":
                data = data.replace(
                    b"presentationml.template.main+xml",
                    b"presentationml.presentation.main+xml",
                )
            contents[item] = data
    with zipfile.ZipFile(tmp.name, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in contents.items():
            zout.writestr(name, data)
    return tmp.name


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=COLOR_BODY, alignment=PP_ALIGN.LEFT,
                font_name="メイリオ"):
    """テキストボックスを追加"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_textbox(slide, left, top, width, height, items, font_size=14,
                       color=COLOR_BODY, font_name="メイリオ", level_indent=None):
    """箇条書きテキストボックスを追加"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # レベルに応じたインデントとマーカー
        level = 0
        text = item
        if isinstance(item, tuple):
            text, level = item

        if level == 0:
            p.text = f"• {text}"
        else:
            p.text = f"  ‐ {text}"
            p.level = 1

        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(4)
    return txBox


# === 16:9 レイアウト定数 (13.333" x 7.5") ===
SLIDE_W = 13.333
MARGIN_L = 1.2   # 左マージン
MARGIN_R = 1.2   # 右マージン
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R  # ≒ 10.933"


def add_section_header(slide, text, top=Inches(0.4)):
    """セクション見出しを追加"""
    # 左のアクセントバー
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(MARGIN_L), top, Inches(0.08), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()

    add_textbox(
        slide, Inches(MARGIN_L + 0.2), top - Inches(0.05), Inches(CONTENT_W - 0.2), Inches(0.6),
        text, font_size=24, bold=True, color=COLOR_HEADING
    )


def add_table(slide, left, top, width, col_widths, headers, rows, font_size=11):
    """テーブルを追加"""
    row_count = len(rows) + 1  # +1 for header
    col_count = len(headers)

    table_shape = slide.shapes.add_table(row_count, col_count, left, top, width, Inches(0.4 * row_count))
    table = table_shape.table

    # 列幅設定
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # ヘッダー行
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = True
            paragraph.font.color.rgb = COLOR_WHITE
            paragraph.font.name = "メイリオ"
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TABLE_HEADER

    # データ行
    for r, row_data in enumerate(rows):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = COLOR_BODY
                paragraph.font.name = "メイリオ"
            # 交互背景
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_TABLE_ALT

    return table_shape


def build_presentation():
    """プレゼンテーションを構築"""
    tmp_pptx = convert_potx_to_pptx(TEMPLATE_PATH)
    prs = Presentation(tmp_pptx)

    # レイアウト参照
    layout_title = prs.slide_layouts[0]    # タイトルスライド
    layout_agenda = prs.slide_layouts[1]   # アジェンダ
    layout_blank = prs.slide_layouts[2]    # 白紙
    layout_closing = prs.slide_layouts[3]  # クロージング

    # =============================================
    # スライド 1: タイトル
    # =============================================
    slide = prs.slides.add_slide(layout_title)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Microsoft 365 Copilot ライセンスで作る\nエージェント"
            for p in ph.text_frame.paragraphs:
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.name = "メイリオ"
        elif ph.placeholder_format.idx == 10:
            ph.text = "Agent Builder と Copilot Studio の違いと使い分け"
            for p in ph.text_frame.paragraphs:
                p.font.size = Pt(16)
                p.font.name = "メイリオ"

    # =============================================
    # スライド 2: アジェンダ
    # =============================================
    slide = prs.slides.add_slide(layout_agenda)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "アジェンダ"
            for p in ph.text_frame.paragraphs:
                p.font.name = "メイリオ"
        elif ph.placeholder_format.idx == 10:
            tf = ph.text_frame
            agenda_items = [
                "Agent Builder とは",
                "Copilot Studio とは",
                "機能比較",
                "どちらを選ぶべきか ― 判断フロー",
                "ガバナンスの違い",
                "Agent Builder → Copilot Studio への移行",
                "ライセンス要件",
                "まとめ",
            ]
            for i, item in enumerate(agenda_items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"{i + 1}. {item}"
                p.font.size = Pt(16)
                p.font.name = "メイリオ"
                p.font.color.rgb = COLOR_BODY
                p.space_after = Pt(6)

    # =============================================
    # スライド 3: はじめに
    # =============================================
    slide = prs.slides.add_slide(layout_blank)
    add_section_header(slide, "はじめに")
    add_textbox(
        slide, Inches(MARGIN_L), Inches(1.2), Inches(CONTENT_W), Inches(1.5),
        "Microsoft 365 Copilot ライセンスを持っていると、Agent Builder（エージェント ビルダー）と "
        "Copilot Studio の 2 つのツールを使ってエージェントを作成できます。",
        font_size=16
    )
    add_textbox(
        slide, Inches(MARGIN_L), Inches(2.5), Inches(CONTENT_W), Inches(1.0),
        "どちらもエージェントを構築できますが、対象ユーザー・機能・ガバナンスの面で大きく異なります。",
        font_size=16
    )

    # 2つのボックスを並べて表示 (16:9 中央配置)
    box_w = 4.8
    box_gap = 0.8
    box_total = box_w * 2 + box_gap
    box_start = (SLIDE_W - box_total) / 2
    box2_start = box_start + box_w + box_gap

    # Agent Builder ボックス
    shape1 = slide.shapes.add_shape(1, Inches(box_start), Inches(3.5), Inches(box_w), Inches(2.5))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
    shape1.line.color.rgb = COLOR_ACCENT
    shape1.line.width = Pt(1.5)
    add_textbox(slide, Inches(box_start + 0.2), Inches(3.6), Inches(box_w - 0.4), Inches(0.5),
                "Agent Builder", font_size=18, bold=True, color=COLOR_ACCENT,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(box_start + 0.2), Inches(4.2), Inches(box_w - 0.4), Inches(1.5),
                "Microsoft 365 Copilot アプリ内に組み込まれた\nエージェント作成機能",
                font_size=13, alignment=PP_ALIGN.CENTER)

    # Copilot Studio ボックス
    shape2 = slide.shapes.add_shape(1, Inches(box2_start), Inches(3.5), Inches(box_w), Inches(2.5))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0xF3, 0xE8, 0xFD)
    shape2.line.color.rgb = RGBColor(0x6B, 0x2F, 0xA0)
    shape2.line.width = Pt(1.5)
    add_textbox(slide, Inches(box2_start + 0.2), Inches(3.6), Inches(box_w - 0.4), Inches(0.5),
                "Copilot Studio", font_size=18, bold=True,
                color=RGBColor(0x6B, 0x2F, 0xA0), alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(box2_start + 0.2), Inches(4.2), Inches(box_w - 0.4), Inches(1.5),
                "スタンドアロンの Web ポータル\n高度なシナリオに対応",
                font_size=13, alignment=PP_ALIGN.CENTER)

    # =============================================
    # スライド 4: Agent Builder とは
    # =============================================
    slide = prs.slides.add_slide(layout_blank)
    add_section_header(slide, "Agent Builder とは")
    add_textbox(
        slide, Inches(MARGIN_L), Inches(1.2), Inches(CONTENT_W), Inches(0.8),
        "Microsoft 365 Copilot アプリ内に組み込まれたエージェント作成機能です。",
        font_size=16
    )
    add_textbox(
        slide, Inches(MARGIN_L), Inches(2.0), Inches(CONTENT_W), Inches(0.5),
        "アクセスポイント：", font_size=14, bold=True, color=COLOR_HEADING
    )
    add_bullet_textbox(
        slide, Inches(MARGIN_L + 0.2), Inches(2.5), Inches(CONTENT_W - 0.2), Inches(1.5),
        [
            "microsoft365.com/chat",
            "office.com/chat",
            "Microsoft Teams（デスクトップ / Web）",
        ],
        font_size=14
    )
    add_textbox(
        slide, Inches(MARGIN_L), Inches(4.0), Inches(CONTENT_W), Inches(0.5),
        "特徴：", font_size=14, bold=True, color=COLOR_HEADING
    )
    add_bullet_textbox(
        slide, Inches(MARGIN_L + 0.2), Inches(4.5), Inches(CONTENT_W - 0.2), Inches(2.0),
        [
            "自然言語でインストラクションを記述",
            "Web サイト、SharePoint ファイル、コネクタをナレッジとして指定",
            "コード不要でエージェントを構築",
        ],
        font_size=14
    )

    # =============================================
    # スライド 5: Copilot Studio とは
    # =============================================
    slide = prs.slides.add_slide(layout_blank)
    add_section_header(slide, "Copilot Studio とは")
    add_textbox(
        slide, Inches(MARGIN_L), Inches(1.2), Inches(CONTENT_W), Inches(0.8),
        "copilotstudio.microsoft.com で利用できるスタンドアロンの Web ポータルです。",
        font_size=16
    )
    add_textbox(
        slide, Inches(MARGIN_L), Inches(2.0), Inches(CONTENT_W), Inches(0.5),
        "対応する高度なシナリオ：", font_size=14, bold=True, color=COLOR_HEADING
    )
    add_bullet_textbox(
        slide, Inches(MARGIN_L + 0.2), Inches(2.5), Inches(CONTENT_W - 0.2), Inches(3.5),
        [
            "マルチステップ ワークフロー",
            "外部 API 連携",
            "カスタム コネクタ",
            "Azure AI サービス統合",
            "自律エージェント機能",
            "ALM（バージョン管理、環境分離、RBAC、テレメトリ）",
        ],
        font_size=14
    )

    # =============================================
    # スライド 6: 機能比較（上半分）
    # =============================================
    slide = prs.slides.add_slide(layout_blank)
    add_section_header(slide, "機能比較")

    headers = ["比較項目", "Agent Builder", "Copilot Studio"]
    rows = [
        ["アクセスポイント", "M365 Copilot アプリ", "copilotstudio.microsoft.com"],
        ["想定ユーザー", "インフォメーション ワーカー\n（IT 部門以外も含む）", "メーカー（市民開発者）\n・プロ開発者"],
        ["エージェントの対象", "自分自身、または\n小規模チーム", "部署全体・組織全体\n・外部顧客"],
        ["エージェントの種類", "組織ナレッジを活用した\n軽量 Q&A エージェント", "マルチステップ ワークフローや\nビジネス システム統合"],
        ["管理・ガバナンス", "M365 管理センター", "Power Platform 管理センター\n（よりきめ細かい制御）"],
    ]

    tbl_w = 11.0
    tbl_left = (SLIDE_W - tbl_w) / 2
    add_table(
        slide,
        left=Inches(tbl_left), top=Inches(1.2),
        width=Inches(tbl_w),
        col_widths=[Inches(2.5), Inches(4.25), Inches(4.25)],
        headers=headers, rows=rows, font_size=11
    )

    # =============================================
    # スライド 7: 機能比較（主な機能 詳細）
    # =============================================
    slide = prs.slides.add_slide(layout_blank)
    add_section_header(slide, "機能比較 ― 主な機能の詳細")

    col_w = 5.2
    col_gap = 0.6
    col1_left = (SLIDE_W - col_w * 2 - col_gap) / 2
    col2_left = col1_left + col_w + col_gap

    # Agent Builder 側
    add_textbox(slide, Inches(col1_left), Inches(1.2), Inches(col_w), Inches(0.4),
                "Agent Builder", font_size=16, bold=True, color=COLOR_ACCENT)
    add_bullet_textbox(
        slide, Inches(col1_left), Inches(1.7), Inches(col_w), Inches(4.5),
        [
            "自然言語によるオーサリング",
            "Microsoft Graph ベースの Q&A",
            "ユーザーの M365 アクセス許可を尊重",
            "M365 Copilot オーケストレーター・\n  基盤モデルを使用",
        ],
        font_size=13
    )

    # Copilot Studio 側
    add_textbox(slide, Inches(col2_left), Inches(1.2), Inches(col_w), Inches(0.4),
                "Copilot Studio", font_size=16, bold=True,
                color=RGBColor(0x6B, 0x2F, 0xA0))
    add_bullet_textbox(
        slide, Inches(col2_left), Inches(1.7), Inches(col_w), Inches(4.5),
        [
            "広範囲への公開（外部含む）",
            "マルチステップ ロジック・承認・\n  分岐ワークフロー",
            "高度な AI モデルと\n  Azure AI サービスの統合",
            "プリビルド / カスタム コネクタ\n  による外部データ接続",
            "自律エージェント機能",
            "ALM（バージョン管理、\n  dev/test/prod 環境、RBAC）",
        ],
        font_size=13
    )

    # =============================================
    # スライド 8: 判断フロー
    # =============================================
    slide = prs.slides.add_slide(layout_blank)
    add_section_header(slide, "どちらを選ぶべきか ― 判断フロー")

    add_textbox(
        slide, Inches(MARGIN_L), Inches(1.2), Inches(CONTENT_W), Inches(0.5),
        "次の 4 つの観点で検討します：", font_size=16
    )

    criteria = [
        ("1", "対象ユーザー", "誰がエージェントを使うか？"),
        ("2", "展開スコープ", "どの範囲に共有するか？"),
        ("3", "機能", "エージェントにどんなタスクを実行させるか？"),
        ("4", "ガバナンスのニーズ", "ALM が必要か？"),
    ]

    criteria_left = 2.0  # 丸の左位置
    for i, (num, title, desc) in enumerate(criteria):
        y = Inches(2.0) + Inches(1.0) * i
        # 番号の丸
        shape = slide.shapes.add_shape(
            9,  # Oval
            Inches(criteria_left), y, Inches(0.5), Inches(0.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_ACCENT
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLOR_WHITE
        tf.paragraphs[0].font.name = "メイリオ"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_textbox(slide, Inches(criteria_left + 0.7), y, Inches(3.0), Inches(0.5),
                    title, font_size=16, bold=True, color=COLOR_HEADING)
        add_textbox(slide, Inches(criteria_left + 3.7), y, Inches(6.5), Inches(0.5),
                    desc, font_size=14)

    # =============================================
    # スライド 9: Agent Builder を選ぶケース
    # =============================================
    slide = prs.slides.add_slide(layout_blank)
    add_section_header(slide, "Agent Builder を選ぶケース")

    add_textbox(slide, Inches(MARGIN_L), Inches(1.2), Inches(CONTENT_W), Inches(0.5),
                "こんなときに最適：", font_size=14, bold=True, color=COLOR_HEADING)
    add_bullet_textbox(
        slide, Inches(MARGIN_L + 0.2), Inches(1.7), Inches(CONTENT_W - 0.2), Inches(1.5),
        [
            "自分やチーム用にすばやくエージェントを作りたい",
            "SharePoint ドキュメントやメールの内容に基づいた Q&A ボットを構築したい",
            "コードは書きたくない（自然言語だけで完結させたい）",
        ],
        font_size=14
    )

    add_textbox(slide, Inches(MARGIN_L), Inches(3.3), Inches(CONTENT_W), Inches(0.5),
                "具体的なユースケース例：", font_size=14, bold=True, color=COLOR_HEADING)

    # ユースケースカード (16:9 中央配置)
    card_w = 3.3
    card_gap = 0.5
    card_total = card_w * 3 + card_gap * 2
    card_start = (SLIDE_W - card_total) / 2
    usecases = [
        ("プロジェクト FAQ ボット", "プロジェクト ドキュメントをもとに\nよくある質問に回答"),
        ("製品ドキュメント\nアシスタント", "社内の製品マニュアルや Wiki から\n情報を検索"),
        ("オンボーディング\nエージェント", "新メンバーが社内ナレッジ ベースから\n回答を得る"),
    ]
    for i, (title, desc) in enumerate(usecases):
        x = Inches(card_start + (card_w + card_gap) * i)
        shape = slide.shapes.add_shape(1, x, Inches(3.9), Inches(card_w), Inches(2.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
        shape.line.color.rgb = COLOR_ACCENT
        shape.line.width = Pt(1)
        add_textbox(slide, x + Inches(0.15), Inches(4.0), Inches(card_w - 0.3), Inches(0.6),
                    title, font_size=13, bold=True, color=COLOR_ACCENT,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), Inches(4.7), Inches(card_w - 0.3), Inches(1.2),
                    desc, font_size=11, alignment=PP_ALIGN.CENTER)

    # =============================================
    # スライド 10: Copilot Studio を選ぶケース
    # =============================================
    slide = prs.slides.add_slide(layout_blank)
    add_section_header(slide, "Copilot Studio を選ぶケース")

    add_textbox(slide, Inches(MARGIN_L), Inches(1.2), Inches(CONTENT_W), Inches(0.5),
                "こんなときに最適：", font_size=14, bold=True, color=COLOR_HEADING)
    add_bullet_textbox(
        slide, Inches(MARGIN_L + 0.2), Inches(1.7), Inches(CONTENT_W - 0.2), Inches(2.0),
        [
            "部署・組織全体、または外部顧客向けにエージェントを展開したい",
            "マルチステップのワークフロー（承認フロー、分岐ロジックなど）が必要",
            "Microsoft 365 以外の外部システムや API と統合したい",
            "エンタープライズ レベルのガバナンス（ALM、DLP、環境分離など）が求められる",
        ],
        font_size=14
    )

    add_textbox(slide, Inches(MARGIN_L), Inches(3.5), Inches(CONTENT_W), Inches(0.5),
                "具体的なユースケース例：", font_size=14, bold=True, color=COLOR_HEADING)

    # ユースケースカード (16:9 中央配置)
    usecases2 = [
        ("カスタマー サポート\nエージェント", "サポート チケットの作成、\n人間へのエスカレーション"),
        ("IT ヘルプデスク\nトリアージ", "IT リクエストを受け付け、\n適切なチームに振り分け"),
        ("営業アシスタント\n（CRM 連携）", "売上データの取得、\nメモ作成、承認ワークフロー"),
    ]
    for i, (title, desc) in enumerate(usecases2):
        x = Inches(card_start + (card_w + card_gap) * i)
        shape = slide.shapes.add_shape(1, x, Inches(4.1), Inches(card_w), Inches(2.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF3, 0xE8, 0xFD)
        shape.line.color.rgb = RGBColor(0x6B, 0x2F, 0xA0)
        shape.line.width = Pt(1)
        add_textbox(slide, x + Inches(0.15), Inches(4.2), Inches(card_w - 0.3), Inches(0.6),
                    title, font_size=13, bold=True,
                    color=RGBColor(0x6B, 0x2F, 0xA0), alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), Inches(4.9), Inches(card_w - 0.3), Inches(1.2),
                    desc, font_size=11, alignment=PP_ALIGN.CENTER)

    # =============================================
    # スライド 11: ガバナンスの違い
    # =============================================
    slide = prs.slides.add_slide(layout_blank)
    add_section_header(slide, "ガバナンスの違い")

    gov_col_w = 5.2
    gov_gap = 0.6
    gov_col1 = (SLIDE_W - gov_col_w * 2 - gov_gap) / 2
    gov_col2 = gov_col1 + gov_col_w + gov_gap

    # Agent Builder 側
    add_textbox(slide, Inches(gov_col1), Inches(1.2), Inches(gov_col_w), Inches(0.4),
                "Agent Builder", font_size=16, bold=True, color=COLOR_ACCENT)
    shape_ab = slide.shapes.add_shape(1, Inches(gov_col1), Inches(1.7), Inches(gov_col_w), Inches(4.5))
    shape_ab.fill.solid()
    shape_ab.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
    shape_ab.line.color.rgb = COLOR_ACCENT
    shape_ab.line.width = Pt(1)
    add_bullet_textbox(
        slide, Inches(gov_col1 + 0.2), Inches(1.9), Inches(gov_col_w - 0.4), Inches(4.0),
        [
            "新しい特権なし ― 既存の M365 アクセス許可を尊重",
            "組み込みの監査 ― 標準監査ログ、DLP/保持ポリシーがそのまま適用",
            "M365 管理センター > Copilot > Agents で管理",
        ],
        font_size=12
    )

    # Copilot Studio 側
    add_textbox(slide, Inches(gov_col2), Inches(1.2), Inches(gov_col_w), Inches(0.4),
                "Copilot Studio", font_size=16, bold=True,
                color=RGBColor(0x6B, 0x2F, 0xA0))
    shape_cs = slide.shapes.add_shape(1, Inches(gov_col2), Inches(1.7), Inches(gov_col_w), Inches(4.5))
    shape_cs.fill.solid()
    shape_cs.fill.fore_color.rgb = RGBColor(0xF3, 0xE8, 0xFD)
    shape_cs.line.color.rgb = RGBColor(0x6B, 0x2F, 0xA0)
    shape_cs.line.width = Pt(1)
    add_bullet_textbox(
        slide, Inches(gov_col2 + 0.2), Inches(1.9), Inches(gov_col_w - 0.4), Inches(4.0),
        [
            "構造化開発（ALM） ― dev / test / prod 環境分離",
            "コネクタ ガバナンス ― 接続先システムを管理者が制御",
            "環境レベルのポリシー ― DLP、RBAC、監査を環境単位で適用",
            "柔軟な展開 ― Teams、Web サイト、カスタム エンドポイントへ公開",
            "Power Platform 管理センターで管理",
        ],
        font_size=12
    )

    # =============================================
    # スライド 12: 移行パス
    # =============================================
    slide = prs.slides.add_slide(layout_blank)
    add_section_header(slide, "Agent Builder → Copilot Studio への移行")

    add_textbox(
        slide, Inches(MARGIN_L), Inches(1.2), Inches(CONTENT_W), Inches(0.8),
        "Agent Builder で作ったエージェントは、後から Copilot Studio にコピーできます。\n"
        "「スモールスタート」のアプローチが有効です。",
        font_size=16
    )

    # ステップの矢印フロー (16:9 中央配置)
    step_box_w = 2.3
    step_arrow_w = 0.5
    step_total = step_box_w * 4 + step_arrow_w * 3
    step_start = (SLIDE_W - step_total) / 2
    steps = [
        "Agent Builder\nで開始",
        "要件が\n複雑化",
        "Copilot Studio\nにコピー",
        "高度な機能\nを追加",
    ]
    for i, step_text in enumerate(steps):
        x = Inches(step_start + (step_box_w + step_arrow_w) * i)
        # ボックス
        shape = slide.shapes.add_shape(
            12,  # Rounded Rectangle
            x, Inches(2.6), Inches(step_box_w), Inches(1.2)
        )
        shape.fill.solid()
        if i < 2:
            shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
            shape.line.color.rgb = COLOR_ACCENT
        else:
            shape.fill.fore_color.rgb = RGBColor(0xF3, 0xE8, 0xFD)
            shape.line.color.rgb = RGBColor(0x6B, 0x2F, 0xA0)
        shape.line.width = Pt(1.5)

        add_textbox(slide, x + Inches(0.1), Inches(2.7), Inches(step_box_w - 0.2), Inches(1.0),
                    step_text, font_size=13, bold=True, color=COLOR_HEADING,
                    alignment=PP_ALIGN.CENTER)

        # 矢印
        if i < len(steps) - 1:
            arrow_x = Inches(step_start + (step_box_w + step_arrow_w) * i + step_box_w)
            add_textbox(slide, arrow_x, Inches(2.9), Inches(step_arrow_w), Inches(0.5),
                       "→", font_size=24, bold=True, color=COLOR_ACCENT,
                       alignment=PP_ALIGN.CENTER)

    add_textbox(
        slide, Inches(MARGIN_L), Inches(4.2), Inches(CONTENT_W), Inches(0.5),
        "コピー時に引き継がれる構成内容：",
        font_size=14, bold=True, color=COLOR_HEADING
    )
    add_bullet_textbox(
        slide, Inches(MARGIN_L + 0.2), Inches(4.7), Inches(CONTENT_W - 0.2), Inches(2.0),
        [
            "エージェントのコア構成とインストラクション",
            "→ Copilot Studio 側で ALM・監視・エンタープライズ ガバナンス・外部統合を追加",
        ],
        font_size=14
    )

    # =============================================
    # スライド 13: ライセンス要件
    # =============================================
    slide = prs.slides.add_slide(layout_blank)
    add_section_header(slide, "ライセンス要件")

    headers_lic = ["条件", "Agent Builder", "Copilot Studio"]
    rows_lic = [
        ["Microsoft 365 Copilot\nアドオン ライセンス", "✅ 利用可能", "✅ 利用可能"],
        ["Copilot Credits /\n従量課金", "✅ 利用可能", "✅ 利用可能"],
        ["無料（指示 + 公開 Web\nのナレッジのみ）", "✅ 利用可能", "❌"],
    ]

    lic_tbl_w = 10.0
    lic_tbl_left = (SLIDE_W - lic_tbl_w) / 2
    add_table(
        slide,
        left=Inches(lic_tbl_left), top=Inches(1.4),
        width=Inches(lic_tbl_w),
        col_widths=[Inches(3.6), Inches(3.2), Inches(3.2)],
        headers=headers_lic, rows=rows_lic, font_size=13
    )

    add_textbox(
        slide, Inches(MARGIN_L), Inches(4.0), Inches(CONTENT_W), Inches(1.0),
        "Microsoft 365 Copilot ライセンスがあれば、どちらのツールも追加費用なしで利用できます。",
        font_size=16, bold=True, color=COLOR_ACCENT
    )

    # =============================================
    # スライド 14: まとめ
    # =============================================
    slide = prs.slides.add_slide(layout_blank)
    add_section_header(slide, "まとめ")

    headers_sum = ["選択基準", "Agent Builder", "Copilot Studio"]
    rows_sum = [
        ["すばやく Q&A ボットを作りたい", "⭐ おすすめ", ""],
        ["コードなしで完結したい", "⭐ おすすめ", ""],
        ["小規模チーム向け", "⭐ おすすめ", ""],
        ["組織全体・外部向けに展開", "", "⭐ おすすめ"],
        ["外部システム連携・ワークフロー", "", "⭐ おすすめ"],
        ["エンタープライズ ガバナンス", "", "⭐ おすすめ"],
    ]

    sum_tbl_w = 11.0
    sum_tbl_left = (SLIDE_W - sum_tbl_w) / 2
    add_table(
        slide,
        left=Inches(sum_tbl_left), top=Inches(1.2),
        width=Inches(sum_tbl_w),
        col_widths=[Inches(4.2), Inches(3.4), Inches(3.4)],
        headers=headers_sum, rows=rows_sum, font_size=12
    )

    # キーメッセージ
    msg_w = 10.5
    msg_left = (SLIDE_W - msg_w) / 2
    shape_msg = slide.shapes.add_shape(
        12, Inches(msg_left), Inches(5.2), Inches(msg_w), Inches(1.0)
    )
    shape_msg.fill.solid()
    shape_msg.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
    shape_msg.line.color.rgb = COLOR_ACCENT
    shape_msg.line.width = Pt(2)
    add_textbox(
        slide, Inches(msg_left + 0.2), Inches(5.3), Inches(msg_w - 0.4), Inches(0.8),
        "迷ったら、まず Agent Builder で小さく始めてみましょう。\n"
        "要件が拡大した時点で Copilot Studio にコピーして拡張できます。",
        font_size=16, bold=True, color=COLOR_ACCENT, alignment=PP_ALIGN.CENTER
    )

    # =============================================
    # スライド 15: クロージング
    # =============================================
    slide = prs.slides.add_slide(layout_closing)

    # === 保存 ===
    prs.save(OUTPUT_PATH)
    os.remove(tmp_pptx)
    print(f"✅ PowerPoint を保存しました: {OUTPUT_PATH}")


if __name__ == "__main__":
    build_presentation()
